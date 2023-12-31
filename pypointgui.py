### Program Coded/Created by KetaGod | keta666 on Discord ###
### For Questions, Comments, or Requests. Please msg me   ###
### via Discord or my email @ ketagod666@proton.me        ###

import tkinter as tk
from tkinter import filedialog
from tkinter import ttk
from pptx import Presentation
from docx import Document
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.text_rank import TextRankSummarizer
from sumy.nlp.stemmers import Stemmer
from sumy.utils import get_stop_words
import subprocess
import sys
import random

# List of available file types
file_types = [
    "Word (docx)",
    "Plain Text (txt)",
    "Markdown (md)",
    "Comma-Separted Values (csv)",
    "JSON (json)",
    "XML (xml)",
    "Rich Text Format (rtf)" 
]

def get_file_extension(selected_type):
    extensions = {
        "Word (docx)": ".docx",
        "Plain Text (txt)": ".txt",
        "Markdown (md)": ".md",
        "Comma-Seperated Values (csv)": ".csv",
        "JSON (json)": ".json",
        "XML (xml)": ".xml",
        "Rich Text Format (rtf)": ".rtf"
    }
    return extensions.get(selected_type, ".txt") # Default output file (feel free to change)

def extract_text_from_pptx(pptx_file, word_output_file):
    selected_type = file_type_combobox.get()
    file_extension = get_file_extension(selected_type)
    if not file_extension:
        return  # Handle the case where no extension is selected

    presentation = Presentation(pptx_file)

    all_text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text += shape.text + "\n"

    # Save to the selected file type
    output_file = word_output_file + file_extension
    if selected_type == "Word (docx)":
        doc = Document()
        doc.add_paragraph(all_text)
        doc.save(output_file)
    elif selected_type == "Plain Text (txt)":
        with open(output_file, "w", encoding="utf-8") as file:
            file.write(all_text)

    status_label.config(text=f"Text extracted and saved to '{output_file}'")

def summarize_text_with_length(text, summary_length=3):
    # Initialize TextRankSummarizer with the specified language
    summarizer = TextRankSummarizer(Stemmer("english"))

    # Set the summarization ratio based on the summary length
    summarizer.stop_words = get_stop_words("english")
    summarized_sentences = summarizer(PlaintextParser.from_string(text, Tokenizer("english")))

    # Select the specified number of sentences for the summary
    num_sentences = min(len(summarized_sentences), summary_length)
    summary = " ".join(str(sentence) for sentence in summarized_sentences[:num_sentences])

    return summary

def summarize_pptx_with_length():
    pptx_file = pptx_entry.get()
    word_output_file = word_output_entry.get()

    presentation = Presentation(pptx_file)

    all_text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                all_text += shape.text + "\n"

    # Prompt the user for the summary length
    summary_length = int(input("Enter the desired summary length (number of sentences): "))

    selected_type = file_type_combobox.get()
    file_extension = get_file_extension(selected_type)
    if not file_extension:
        return  # Handle the case where no extension is selected

    # Summarize the text with the specified length
    summarized_text = summarize_text_with_length(all_text, summary_length)

    # Save the summary to the selected file type
    output_file = word_output_file + file_extension
    if selected_type == "Word (docx)":
        doc = Document()
        doc.add_paragraph(summarized_text)
        doc.save(output_file)
    elif selected_type == "Plain Text (txt)":
        with open(output_file, "w", encoding="utf-8") as file:
            file.write(summarized_text)

    status_label.config(text=f"Text summarized and saved to '{output_file}'")


def browse_pptx():
    pptx_file = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    pptx_entry.delete(0, "end")
    pptx_entry.insert(0, pptx_file)

def browse_word_output():
    word_output_file = filedialog.asksaveasfilename(defaultextension=".docx", filetypes=[("Word files", "*.docx")])
    word_output_entry.delete(0, "end")
    word_output_entry.insert(0, word_output_file)

def extract_text():
    pptx_file = pptx_entry.get()
    word_output_file = word_output_entry.get()
    extract_text_from_pptx(pptx_file, word_output_file)
    status_label.config(text=f"Text extracted and saved to '{word_output_file}'")

def run_script_silently():
    script = "pypointgui.py"
    startupinfo = subprocess.STARTUPINFO()
    startupinfo.dwFlags |= subprocess.STARTF_USESHOWWINDOW
    subprocess.Popen([sys.executable, script], startupinfo=startupinfo)

# Create a Tkinter window
window = tk.Tk()
window.title("PowerPoint Notes ExtractorV0.0.4")

# Allow resizing of the GUI
window.geometry("450x270") # Initial Size

# Set a random background color
random_color = "#{:02x}{:02x}{:02x}".format(random.randint(0, 255), random.randint(0, 255), random.randint(0, 255))
window.configure(bg=random_color)

# Apply high contrast theme for better clarity
window.tk_setPalette(background=random_color, foreground='white')

# Add a button to run text summarization with user-defined length
summarize_with_length_button = tk.Button(window, text="Summarize Text (Custom Length)", command=summarize_pptx_with_length)
summarize_with_length_button.grid(row=8, column=0, columnspan=3)

# Create and configure GUI elements
author_label = tk.Label(window, text="Author: KetaGod | keta666")
author_label.grid(row=0, column=0, columnspan=3)

pptx_label = tk.Label(window, text="Select PowerPoint file:")
pptx_entry = tk.Entry(window, width=40)
browse_pptx_button = tk.Button(window, text="Browse", command=browse_pptx)

word_output_label = tk.Label(window, text="Select output Word file:")
word_output_entry = tk.Entry(window, width=40)
browse_word_output_button = tk.Button(window, text="Browse", command=browse_word_output)

extract_button = tk.Button(window, text="Extract Text", command=extract_text)
status_label = tk.Label(window, text="")

word_output_label.grid(row=2, column=0)
word_output_entry.grid(row=2, column=1)
browse_word_output_button.grid(row=2, column=2)

extract_button.grid(row=3, column=0, columnspan=3)
status_label.grid(row=4, column=0, columnspan=3)

# Add combobox for selceting file types
file_type_label = tk.Label(window, text="Select Output File Type:")
file_type_combobox = ttk.Combobox(window, values=file_types)
file_type_combobox.set(file_types[0]) # Sets base selection (Free to change)
file_type_label.grid(row=5, column=0)
file_type_combobox.grid(row=5, column=1)

# Add a multi-line information label
info_text = (
    "Thank you for using PyPoint.\n"
    "If you haven't already, please read the README.txt file.\n"
    "This program was created for those who do not want to take their own notes on PowerPoints.\n"
    "This program is able to write everything down within a second or less.\n"
    "Enjoy. Created by KetaGod"
)
info_label = tk.Label(window, text=info_text, bg=random_color, wraplength=380, justify="left")
info_label.grid(row=6, column=0, columnspan=3)

# Arrange GUI elements
pptx_label.grid(row=1, column=0)
pptx_entry.grid(row=1, column=1)
browse_pptx_button.grid(row=1, column=2)

word_output_label.grid(row=2, column=0)
word_output_entry.grid(row=2, column=1)
browse_word_output_button.grid(row=2, column=2)

extract_button.grid(row=3, column=0, columnspan=3)
status_label.grid(row=4, column=0, columnspan=3)

# Start the GUI main loop
window.mainloop()
